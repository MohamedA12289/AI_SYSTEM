
from __future__ import annotations

from pathlib import Path
from datetime import datetime, timezone
import hashlib
import json
import mimetypes
import shutil
import subprocess
import uuid
import zipfile

import pandas as pd
try:
    import pytesseract
    _PYTESSERACT_AVAILABLE = True
except ImportError:
    pytesseract = None  # type: ignore[assignment]
    _PYTESSERACT_AVAILABLE = False
from PIL import Image
from pypdf import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation

from memory import validate_project_name
from project_registry import create_project
from ingest_store import (
    ensure_project_ingest_dirs,
    create_job,
    update_job,
    upsert_document,
    list_documents,
    get_document,
    write_text_artifact,
    get_project_ingest_root,
)
from process_utils import run_hidden

try:
    from tika import parser as tika_parser  # type: ignore
except Exception:
    tika_parser = None

MAX_INDEXED_FILES_PER_JOB = 4000
MAX_INLINE_TEXT_CHARS = 30000
MAX_TIKA_FILE_SIZE_BYTES = 250 * 1024 * 1024
MAX_DIRECT_TEXT_READ_BYTES = 8 * 1024 * 1024
MAX_PDF_PAGES = 20
TABULAR_SAMPLE_ROWS = 200

TEXT_EXTENSIONS = {
    ".txt", ".md", ".rst", ".log", ".json", ".yaml", ".yml", ".toml", ".xml", ".csv", ".tsv",
    ".ini", ".env", ".sql", ".py", ".js", ".ts", ".tsx", ".jsx", ".java", ".go", ".rs", ".cpp",
    ".c", ".h", ".hpp", ".cs", ".php", ".rb", ".swift", ".kt", ".html", ".css", ".scss", ".less",
    ".sh", ".bat", ".ps1",
}
PDF_EXTENSIONS = {".pdf"}
DOCX_EXTENSIONS = {".docx"}
SLIDE_EXTENSIONS = {".pptx"}
SPREADSHEET_EXTENSIONS = {".csv", ".tsv", ".xlsx", ".xls"}
IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".webp", ".gif", ".bmp", ".tiff", ".svg"}
AUDIO_EXTENSIONS = {".mp3", ".wav", ".m4a", ".flac"}
VIDEO_EXTENSIONS = {".mp4", ".mov", ".mkv", ".avi", ".webm"}
ARCHIVE_EXTENSIONS = {".zip", ".7z", ".tar", ".gz", ".bz2"}


def _now_iso() -> str:
    return datetime.now(timezone.utc).isoformat()


def supported_file_families() -> dict:
    return {
        "text_code_config": sorted(TEXT_EXTENSIONS),
        "documents": [".pdf", ".docx", ".pptx"],
        "spreadsheets": sorted(SPREADSHEET_EXTENSIONS),
        "images": sorted(IMAGE_EXTENSIONS),
        "audio": sorted(AUDIO_EXTENSIONS),
        "video": sorted(VIDEO_EXTENSIONS),
        "archives": sorted(ARCHIVE_EXTENSIONS),
        "special": ["folder"],
    }


def validate_local_source_path(source_path: str) -> Path:
    if not isinstance(source_path, str):
        raise ValueError("source_path must be a string.")

    cleaned = source_path.strip()
    if not cleaned:
        raise ValueError("source_path cannot be empty.")

    path = Path(cleaned)
    if not path.is_absolute():
        raise ValueError("source_path must be an absolute local path.")

    if not path.exists():
        raise FileNotFoundError("source_path does not exist.")

    return path.resolve()


def infer_source_kind(source_path: Path) -> str:
    if source_path.is_dir():
        return "folder"
    if source_path.suffix.lower() == ".zip":
        return "zip"
    return "file"


def _safe_name(value: str) -> str:
    cleaned = "".join(ch if ch.isalnum() or ch in ("-", "_", ".") else "_" for ch in value)
    return cleaned[:120] or "item"


def _sha256_file(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        while True:
            chunk = handle.read(1024 * 1024)
            if not chunk:
                break
            digest.update(chunk)
    return digest.hexdigest()


def _read_text_head(path: Path, limit_bytes: int = MAX_DIRECT_TEXT_READ_BYTES) -> str:
    with path.open("rb") as handle:
        raw = handle.read(limit_bytes)
    return raw.decode("utf-8", errors="ignore")[:MAX_INLINE_TEXT_CHARS]


def _extract_pdf_text(path: Path) -> str:
    parts: list[str] = []
    reader = PdfReader(str(path))
    for page in reader.pages[:MAX_PDF_PAGES]:
        try:
            parts.append(page.extract_text() or "")
        except Exception:
            continue
        if sum(len(p) for p in parts) >= MAX_INLINE_TEXT_CHARS:
            break
    return "\n".join(parts)[:MAX_INLINE_TEXT_CHARS]


def _extract_docx_text(path: Path) -> str:
    doc = DocxDocument(str(path))
    text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    return text[:MAX_INLINE_TEXT_CHARS]


def _extract_pptx_text(path: Path) -> str:
    prs = Presentation(str(path))
    parts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if text:
                parts.append(text)
        if sum(len(p) for p in parts) >= MAX_INLINE_TEXT_CHARS:
            break
    return "\n".join(parts)[:MAX_INLINE_TEXT_CHARS]


def _extract_image_text(path: Path) -> str:
    if not _PYTESSERACT_AVAILABLE:
        return ""
    try:
        with Image.open(path) as image:
            text = pytesseract.image_to_string(image)
        return text[:MAX_INLINE_TEXT_CHARS]
    except Exception:
        return ""


def _ffprobe_metadata(path: Path) -> dict:
    command = [
        "ffprobe",
        "-v", "error",
        "-print_format", "json",
        "-show_format",
        "-show_streams",
        str(path),
    ]
    result = run_hidden(command, capture_output=True, text=True, timeout=30)
    if result.returncode != 0:
        return {
            "probe_ok": False,
            "stderr": (result.stderr or "").strip()[:1000],
        }
    try:
        data = json.loads(result.stdout or "{}")
    except Exception:
        data = {"raw": (result.stdout or "")[:2000]}
    data["probe_ok"] = True
    return data


def _tika_parse(path: Path) -> tuple[str, dict]:
    if tika_parser is None:
        return "", {"tika_available": False}

    if path.stat().st_size > MAX_TIKA_FILE_SIZE_BYTES:
        return "", {"tika_available": True, "skipped_reason": "file_too_large_for_tika"}

    try:
        parsed = tika_parser.from_file(str(path))
        metadata = parsed.get("metadata") or {}
        content = parsed.get("content") or ""
        return str(content)[:MAX_INLINE_TEXT_CHARS], metadata if isinstance(metadata, dict) else {}
    except Exception as exc:
        return "", {"tika_available": True, "error": str(exc)[:500]}


def classify_extension(path: Path) -> tuple[str, str]:
    ext = path.suffix.lower()
    if ext in TEXT_EXTENSIONS:
        return "text", ext
    if ext in PDF_EXTENSIONS or ext in DOCX_EXTENSIONS or ext in SLIDE_EXTENSIONS:
        return "document", ext
    if ext in SPREADSHEET_EXTENSIONS:
        return "tabular", ext
    if ext in IMAGE_EXTENSIONS:
        return "image", ext
    if ext in AUDIO_EXTENSIONS:
        return "audio", ext
    if ext in VIDEO_EXTENSIONS:
        return "video", ext
    if ext in ARCHIVE_EXTENSIONS:
        return "archive", ext
    return "unknown", ext


def summarize_tabular_path(path: Path) -> dict:
    ext = path.suffix.lower()

    if ext == ".csv":
        df = pd.read_csv(path, nrows=TABULAR_SAMPLE_ROWS)
        sheet_names = None
    elif ext == ".tsv":
        df = pd.read_csv(path, sep="\t", nrows=TABULAR_SAMPLE_ROWS)
        sheet_names = None
    elif ext in {".xlsx", ".xls"}:
        excel = pd.ExcelFile(path)
        sheet_names = excel.sheet_names
        first_sheet = sheet_names[0] if sheet_names else 0
        df = pd.read_excel(path, sheet_name=first_sheet, nrows=TABULAR_SAMPLE_ROWS)
    else:
        raise ValueError("Document is not a supported tabular file.")

    preview_rows = df.head(10).fillna("").to_dict(orient="records")
    dtypes = {str(col): str(dtype) for col, dtype in df.dtypes.items()}

    numeric_summary = {}
    numeric_df = df.select_dtypes(include=["number"])
    if not numeric_df.empty:
        describe = numeric_df.describe().fillna("").to_dict()
        for column, stats in describe.items():
            numeric_summary[str(column)] = {str(k): v for k, v in stats.items()}

    return {
        "columns": [str(col) for col in df.columns.tolist()],
        "row_count_sampled": int(len(df)),
        "preview_rows": preview_rows,
        "dtypes": dtypes,
        "numeric_summary": numeric_summary,
        "sheet_names": sheet_names,
    }


def _extract_primary_text(path: Path, family: str, ext: str) -> tuple[str, str, dict]:
    extra: dict = {}
    parser_name = "metadata_only"

    if family == "text":
        parser_name = "direct_text"
        return _read_text_head(path), parser_name, extra

    if ext == ".pdf":
        try:
            return _extract_pdf_text(path), "pypdf", extra
        except Exception:
            text, metadata = _tika_parse(path)
            return text, "tika_pdf", {"tika_metadata": metadata}

    if ext == ".docx":
        try:
            return _extract_docx_text(path), "python_docx", extra
        except Exception:
            text, metadata = _tika_parse(path)
            return text, "tika_docx", {"tika_metadata": metadata}

    if ext == ".pptx":
        try:
            return _extract_pptx_text(path), "python_pptx", extra
        except Exception:
            text, metadata = _tika_parse(path)
            return text, "tika_pptx", {"tika_metadata": metadata}

    if family == "tabular":
        try:
            summary = summarize_tabular_path(path)
            text = json.dumps({
                "columns": summary.get("columns"),
                "row_count_sampled": summary.get("row_count_sampled"),
                "sheet_names": summary.get("sheet_names"),
                "preview_rows": summary.get("preview_rows", [])[:5],
            }, indent=2)
            return text[:MAX_INLINE_TEXT_CHARS], "pandas", extra
        except Exception as exc:
            extra["tabular_error"] = str(exc)[:500]
            text, metadata = _tika_parse(path)
            return text, "tika_tabular", {"tika_metadata": metadata, **extra}

    if family == "image":
        text = _extract_image_text(path)
        try:
            with Image.open(path) as image:
                extra["image_width"] = image.width
                extra["image_height"] = image.height
                extra["image_mode"] = image.mode
        except Exception:
            pass
        return text, "tesseract_ocr", extra

    if family in {"audio", "video"}:
        extra["media_metadata"] = _ffprobe_metadata(path)
        return "", "ffprobe", extra

    if family == "archive":
        if ext == ".zip":
            try:
                with zipfile.ZipFile(path, "r") as archive:
                    names = archive.namelist()[:200]
                text = "\n".join(names)[:MAX_INLINE_TEXT_CHARS]
                return text, "zip_manifest", extra
            except Exception:
                pass
        text, metadata = _tika_parse(path)
        return text, "tika_archive", {"tika_metadata": metadata}

    text, metadata = _tika_parse(path)
    if text:
        return text, "tika_fallback", {"tika_metadata": metadata}

    return "", parser_name, extra


def search_documents(project_name: str, query: str, limit: int = 20) -> dict:
    cleaned = (query or "").strip().lower()
    if not cleaned:
        raise ValueError("query cannot be empty.")

    results = []
    for document in list_documents(project_name):
        haystack_parts = [
            str(document.get("file_name", "")),
            str(document.get("relative_path", "")),
            str(document.get("file_family", "")),
            str(document.get("text_excerpt", "")),
            str(document.get("summary", "")),
        ]
        haystack = "\n".join(haystack_parts).lower()
        score = 0
        for token in cleaned.split():
            if token in haystack:
                score += 1
        if score > 0:
            results.append({
                "document_id": document.get("document_id"),
                "file_name": document.get("file_name"),
                "relative_path": document.get("relative_path"),
                "file_family": document.get("file_family"),
                "score": score,
                "excerpt": str(document.get("text_excerpt", ""))[:600],
            })

    results.sort(key=lambda item: (-item["score"], item["file_name"] or ""))
    return {
        "query": query,
        "results": results[:limit],
        "total_hits": len(results),
    }


def summarize_document(project_name: str, document_id: str) -> dict:
    document = get_document(project_name, document_id)
    summary = document.get("summary")
    if summary:
        return {"document_id": document_id, "summary": summary, "cached": True}

    text_path_rel = document.get("text_artifact_path")
    content = document.get("text_excerpt", "")
    if text_path_rel:
        full_path = get_project_ingest_root(project_name) / text_path_rel
        if full_path.exists():
            content = full_path.read_text(encoding="utf-8")

    if not content.strip():
        summary = "No extractable text was available. Use metadata, file structure, or specialized tools for this file."
    else:
        prompt = f"""
Summarize this file for a coding/project assistant.
Be concise and practical.
Include:
- what this file appears to be
- the key useful points
- anything worth following up on

File name: {document.get('file_name')}
File kind: {document.get('file_family')}
Text/content:
{content[:12000]}
"""
        from ollama_client import ask_ollama  # lazy import to avoid circular dependency
        summary = ask_ollama(prompt)

    document = dict(document)
    document["summary"] = summary
    document["summary_updated_at"] = _now_iso()
    upsert_document(project_name, document)

    return {"document_id": document_id, "summary": summary, "cached": False}


def _iter_files(root: Path) -> list[Path]:
    files: list[Path] = []
    for path in root.rglob("*"):
        if path.is_file():
            files.append(path)
        if len(files) >= MAX_INDEXED_FILES_PER_JOB:
            break
    return files


def _copy_file(src: Path, dst_root: Path) -> Path:
    dst_root.mkdir(parents=True, exist_ok=True)
    target = dst_root / src.name
    if target.exists():
        target = dst_root / f"{uuid.uuid4().hex[:8]}_{src.name}"
    shutil.copy2(src, target)
    return target


def _copy_folder(src: Path, dst_root: Path) -> Path:
    dst_root.mkdir(parents=True, exist_ok=True)
    target = dst_root / _safe_name(src.name)
    if target.exists():
        target = dst_root / f"{uuid.uuid4().hex[:8]}_{_safe_name(src.name)}"
    shutil.copytree(src, target)
    return target


def _copy_and_extract_zip(src: Path, dst_root: Path) -> Path:
    dst_root.mkdir(parents=True, exist_ok=True)
    archive_dir = dst_root / f"{uuid.uuid4().hex[:8]}_{_safe_name(src.stem)}"
    archive_dir.mkdir(parents=True, exist_ok=True)
    copied_zip = archive_dir / src.name
    shutil.copy2(src, copied_zip)
    extracted_dir = archive_dir / "extracted"
    extracted_dir.mkdir(parents=True, exist_ok=True)
    with zipfile.ZipFile(copied_zip, "r") as archive:
        archive.extractall(extracted_dir)
    return extracted_dir


def _relative_to_root(file_path: Path, root: Path) -> str:
    try:
        return str(file_path.relative_to(root)).replace("\\", "/")
    except Exception:
        return file_path.name


def _build_document_record(project_name: str, file_path: Path, scan_root: Path, access_mode: str, managed_root: str) -> dict:
    stat = file_path.stat()
    family, ext = classify_extension(file_path)
    mime_type, _ = mimetypes.guess_type(str(file_path))
    text_excerpt, parser_name, extra = _extract_primary_text(file_path, family, ext)
    document_id = uuid.uuid4().hex[:12]
    text_artifact_path = write_text_artifact(project_name, document_id, text_excerpt)

    return {
        "document_id": document_id,
        "file_name": file_path.name,
        "relative_path": _relative_to_root(file_path, scan_root),
        "absolute_path": str(file_path),
        "managed_root": managed_root,
        "access_mode": access_mode,
        "extension": ext,
        "mime_type": mime_type,
        "size_bytes": stat.st_size,
        "sha256": _sha256_file(file_path) if stat.st_size <= 512 * 1024 * 1024 else None,
        "file_family": family,
        "parser": parser_name,
        "text_excerpt": (text_excerpt or "")[:3000],
        "text_artifact_path": text_artifact_path,
        "summary": None,
        "indexed_at": _now_iso(),
        **extra,
    }


def ingest_source(project_name: str, source_path: str, source_kind: str | None = None, access_mode: str = "import") -> dict:
    validate_project_name(project_name)
    source = validate_local_source_path(source_path)
    detected_kind = infer_source_kind(source)
    requested_kind = (source_kind or detected_kind).strip().lower()

    if requested_kind not in {"file", "folder", "zip"}:
        raise ValueError("source_kind must be one of: file, folder, zip.")

    if access_mode not in {"import", "link_readonly", "link_readwrite"}:
        raise ValueError("access_mode must be one of: import, link_readonly, link_readwrite.")

    ensure_project_ingest_dirs(project_name)
    job = create_job(project_name, "ingest", str(source), requested_kind, access_mode)
    update_job(project_name, job["job_id"], status="running")

    try:
        sources_root = ensure_project_ingest_dirs(project_name)["sources"] / job["job_id"]
        if access_mode == "import":
            if requested_kind == "file":
                scan_root = _copy_file(source, sources_root)
            elif requested_kind == "folder":
                scan_root = _copy_folder(source, sources_root)
            else:
                scan_root = _copy_and_extract_zip(source, sources_root)
        else:
            scan_root = source

        files_to_index = [scan_root] if scan_root.is_file() else _iter_files(scan_root)
        documents = []
        managed_root = str(scan_root)

        for file_path in files_to_index:
            try:
                record = _build_document_record(
                    project_name=project_name,
                    file_path=file_path,
                    scan_root=scan_root if scan_root.is_dir() else file_path.parent,
                    access_mode=access_mode,
                    managed_root=managed_root,
                )
                upsert_document(project_name, record)
                documents.append(record)
            except Exception as file_error:
                error_record = {
                    "document_id": uuid.uuid4().hex[:12],
                    "file_name": file_path.name,
                    "relative_path": _relative_to_root(file_path, scan_root if scan_root.is_dir() else file_path.parent),
                    "absolute_path": str(file_path),
                    "managed_root": managed_root,
                    "access_mode": access_mode,
                    "extension": file_path.suffix.lower(),
                    "mime_type": None,
                    "size_bytes": file_path.stat().st_size if file_path.exists() else None,
                    "sha256": None,
                    "file_family": "error",
                    "parser": "error",
                    "text_excerpt": "",
                    "text_artifact_path": None,
                    "summary": None,
                    "indexed_at": _now_iso(),
                    "error": str(file_error)[:500],
                }
                upsert_document(project_name, error_record)
                documents.append(error_record)

        final_job = update_job(
            project_name,
            job["job_id"],
            status="completed",
            managed_root=str(scan_root),
            documents_indexed=len(documents),
        )

        return {
            "job": final_job,
            "documents_indexed": len(documents),
            "managed_root": str(scan_root),
            "access_mode": access_mode,
            "source_kind": requested_kind,
        }
    except Exception as exc:
        failed_job = update_job(project_name, job["job_id"], status="error", error=str(exc)[:1000])
        return {
            "job": failed_job,
            "documents_indexed": 0,
            "managed_root": None,
            "access_mode": access_mode,
            "source_kind": requested_kind,
        }


def import_existing_project(project_name: str, display_name: str | None, description: str, source_path: str, access_mode: str = "link_readonly") -> dict:
    try:
        project = create_project(
            project_name=project_name,
            display_name=display_name,
            description=description,
        )
    except FileExistsError:
        project = {
            "project_name": project_name,
            "display_name": display_name or project_name,
            "description": description,
        }

    source = validate_local_source_path(source_path)
    source_kind = infer_source_kind(source)
    if source_kind == "file" and source.suffix.lower() != ".zip":
        source_kind = "folder"

    ingest_result = ingest_source(
        project_name=project_name,
        source_path=str(source),
        source_kind=source_kind,
        access_mode=access_mode,
    )

    return {
        "project": project,
        "ingest": ingest_result,
    }

